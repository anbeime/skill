#!/usr/bin/env python3
"""
生成 PPT 演示文稿
"""
import argparse
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor


def generate_ppt(txt_path: str, output_path: str) -> None:
    """
    根据论文文本生成 PPT 演示文稿

    Args:
        txt_path: 文本文件路径
        output_path: 输出 pptx 文件路径
    """
    print(f"正在生成 PPT: {txt_path}")

    # 读取文本
    with open(txt_path, 'r', encoding='utf-8') as f:
        text = f.read()

    # 创建演示文稿
    prs = Presentation()

    # 设置幻灯片大小（16:9）
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 1. 标题页
    slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    # 尝试提取标题
    title_lines = text.split('\n')
    paper_title = title_lines[0][:50] if title_lines else "论文分析报告"

    title.text = paper_title
    subtitle.text = "自动生成的分析报告"

    # 2. 摘要页
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "摘要"

    # 提取摘要（前 500 个字符）
    content_box = slide.placeholders[1].text_frame
    content_box.text = text[:500] + ("..." if len(text) > 500 else "")

    # 3. 关键点页
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "关键内容"

    # 提取句子
    sentences = re.split(r'[.!?]+', text[:2000])
    sentences = [s.strip() for s in sentences if s.strip() and len(s) > 20][:5]

    content_box = slide.placeholders[1].text_frame
    for i, sentence in enumerate(sentences):
        if i == 0:
            content_box.text = sentence
        else:
            p = content_box.add_paragraph()
            p.text = sentence
            p.level = 0
            p.font.size = Pt(14)

    # 4. 统计信息页
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "统计信息"

    content_box = slide.placeholders[1].text_frame
    stats = [
        f"总字符数: {len(text):,}",
        f"总词数: {len(text.split()):,}",
        f"段落数: {len([p for p in text.split('\\n') if p.strip()]):,}",
        f"句子数: {len(re.split(r'[.!?]+', text)):,}",
    ]

    for i, stat in enumerate(stats):
        if i == 0:
            content_box.text = stat
        else:
            p = content_box.add_paragraph()
            p.text = stat
            p.level = 0

    # 5. 内容页（逐段展示）
    paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
    max_content_slides = 5

    for i in range(min(max_content_slides, len(paragraphs))):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = f"内容详情 {i + 1}"

        content_box = slide.placeholders[1].text_frame
        content_box.text = paragraphs[i][:800] + ("..." if len(paragraphs[i]) > 800 else "")

    # 6. 结论页
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "总结"

    # 提取最后部分作为总结
    content_box = slide.placeholders[1].text_frame
    if len(paragraphs) > 0:
        summary_text = paragraphs[-1][:800] + ("..." if len(paragraphs[-1]) > 800 else "")
        content_box.text = summary_text

    # 确保输出目录存在
    output_dir = os.path.dirname(output_path)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)

    # 保存 PPT
    prs.save(output_path)

    print(f"PPT 已保存到: {output_path}")
    print(f"共生成 {len(prs.slides)} 页幻灯片")


def main():
    parser = argparse.ArgumentParser(description='生成 PPT 演示文稿')
    parser.add_argument('--txt', required=True, help='文本文件路径')
    parser.add_argument('--output', required=True, help='输出 pptx 文件路径')
    args = parser.parse_args()

    generate_ppt(args.txt, args.output)


if __name__ == '__main__':
    main()
