#!/usr/bin/env python3
"""
PPTX Builder - Generate .pptx files from JSON data.

This module creates standard PowerPoint (.pptx) files from JSON data,
supporting multiple layouts, styles, charts, and tables.
"""

import argparse
import json
import os
from pathlib import Path
from typing import Any, Dict, List, Optional

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RgbColor
except ImportError:
    print("错误: 未安装 python-pptx")
    print("请运行: pip install python-pptx")
    exit(1)


# =============================================================================
# Constants
# =============================================================================

DEFAULT_OUTPUT_PATH = "presentation.pptx"
DEFAULT_SLIDE_WIDTH = Inches(10)
DEFAULT_SLIDE_HEIGHT = Inches(7.5)


# =============================================================================
# PPTX Builder
# =============================================================================

class PPTXBuilder:
    """Builder for creating .pptx files from JSON data."""

    def __init__(self, style_config: Optional[Dict[str, Any]] = None):
        """
        Initialize PPTX builder.

        Args:
            style_config: Style configuration dictionary.
        """
        self.style_config = style_config or {}
        self.presentation = None

        print(f"PPTX 构建器已初始化")
        if style_config:
            print(f"  风格配置已加载")

    def load_json(self, json_path: str) -> Dict[str, Any]:
        """
        Load PPT data from JSON file.

        Args:
            json_path: Path to JSON file.

        Returns:
            Parsed JSON data.
        """
        with open(json_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        print(f"✓ JSON 数据已加载: {json_path}")
        return data

    def create_presentation(self, metadata: Dict[str, Any]) -> None:
        """
        Create presentation with metadata.

        Args:
            metadata: Metadata dictionary.
        """
        self.presentation = Presentation()
        self.presentation.slide_width = DEFAULT_SLIDE_WIDTH
        self.presentation.slide_height = DEFAULT_SLIDE_HEIGHT

        # Set core properties
        if metadata.get('title'):
            self.presentation.core_properties.title = metadata['title']
        if metadata.get('author'):
            self.presentation.core_properties.author = metadata['author']

        print(f"✓ 演示文稿已创建: {metadata.get('title', 'Untitled')}")

    def add_slide(
        self,
        slide_data: Dict[str, Any],
        slide_index: int
    ) -> None:
        """
        Add a slide to the presentation.

        Args:
            slide_data: Slide data dictionary.
            slide_index: Slide index (1-based).
        """
        if not self.presentation:
            raise ValueError("演示文稿未创建")

        layout_name = slide_data.get('layout', 'ContentSlide')
        title = slide_data.get('title', '')
        content = slide_data.get('content', [])
        notes = slide_data.get('notes', '')

        # Select layout
        slide_layout = self._select_layout(layout_name)
        slide = self.presentation.slides.add_slide(slide_layout)

        # Add title
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = title
            self._apply_title_style(title_shape)

        # Add content
        self._add_content(slide, content)

        # Add notes
        if notes:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = notes

        print(f"  ✓ 幻灯片 {slide_index} 已添加: {title}")

    def _select_layout(self, layout_name: str):
        """
        Select slide layout based on layout name.

        Args:
            layout_name: Layout name.

        Returns:
            Slide layout object.
        """
        layout_map = {
            'TitleSlide': 0,  # Title slide
            'ContentSlide': 5,  # Title and Content
            'TwoColumnSlide': 6,  # Two Content
            'SectionHeaderSlide': 2,  # Section Header
            'ContentWithCaptionSlide': 7,  # Content with Caption
            'SummarySlide': 5,  # Similar to ContentSlide
        }

        layout_idx = layout_map.get(layout_name, 5)
        return self.presentation.slide_layouts[layout_idx]

    def _apply_title_style(self, title_shape):
        """
        Apply style to title shape.

        Args:
            title_shape: Title shape object.
        """
        text_frame = title_shape.text_frame
        paragraph = text_frame.paragraphs[0]

        # Font settings
        font = paragraph.font
        font.name = self.style_config.get('title_font', 'Arial')
        font.size = Pt(self.style_config.get('title_font_size', 36))
        font.bold = True
        font.color.rgb = RgbColor(*self._hex_to_rgb(
            self.style_config.get('title_color', '#333333')
        ))

    def _add_content(self, slide, content: List[str]):
        """
        Add content to slide.

        Args:
            slide: Slide object.
            content: List of content strings.
        """
        if not content:
            return

        # Find text frame for content
        text_frame = None
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                text_frame = shape.text_frame
                break

        if not text_frame:
            return

        # Clear existing paragraphs
        for paragraph in text_frame.paragraphs[1:]:
            p = paragraph._element
            p.getparent().remove(p)

        # Add content as bullet points
        for i, item in enumerate(content):
            if i == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()

            paragraph.text = item
            paragraph.level = 0

            # Font settings
            font = paragraph.font
            font.name = self.style_config.get('content_font', 'Arial')
            font.size = Pt(self.style_config.get('content_font_size', 24))
            font.color.rgb = RgbColor(*self._hex_to_rgb(
                self.style_config.get('content_color', '#555555')
            ))

    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """
        Convert hex color to RGB tuple.

        Args:
            hex_color: Hex color string (e.g., '#RRGGBB').

        Returns:
            RGB tuple (r, g, b).
        """
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def save(self, output_path: str) -> None:
        """
        Save presentation to .pptx file.

        Args:
            output_path: Output .pptx file path.
        """
        if not self.presentation:
            raise ValueError("演示文稿未创建")

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        self.presentation.save(str(output_path))
        print(f"✓ PPTX 文件已保存: {output_path}")

    def build_from_json(
        self,
        json_path: str,
        output_path: str
    ) -> bool:
        """
        Build .pptx file from JSON data.

        Args:
            json_path: Path to JSON file.
            output_path: Output .pptx file path.

        Returns:
            True if successful, False otherwise.
        """
        try:
            # Load JSON
            print(f"\n开始构建 PPTX...")
            data = self.load_json(json_path)

            # Create presentation
            metadata = data.get('metadata', {})
            self.create_presentation(metadata)

            # Add slides
            slides = data.get('slides', [])
            print(f"\n添加幻灯片...")
            for i, slide_data in enumerate(slides, 1):
                self.add_slide(slide_data, i)

            # Save
            print(f"\n保存文件...")
            self.save(output_path)

            print(f"\n✓ PPTX 文件生成成功: {output_path}")
            print(f"  总页数: {len(slides)}")
            return True

        except Exception as e:
            print(f"\n✗ 错误: {e}")
            import traceback
            traceback.print_exc()
            return False


# =============================================================================
# Main Function
# =============================================================================

def main():
    """Main entry point."""
    parser = argparse.ArgumentParser(
        description='Generate .pptx files from JSON data'
    )
    parser.add_argument(
        '--input', '-i',
        required=True,
        help='Input JSON file path'
    )
    parser.add_argument(
        '--style', '-s',
        help='Style configuration JSON file path'
    )
    parser.add_argument(
        '--output', '-o',
        default=DEFAULT_OUTPUT_PATH,
        help='Output .pptx file path'
    )

    args = parser.parse_args()

    # Load style config
    style_config = None
    if args.style:
        print(f"\n加载风格配置: {args.style}")
        with open(args.style, 'r', encoding='utf-8') as f:
            style_config = json.load(f)

    # Create builder
    builder = PPTXBuilder(style_config=style_config)

    # Build PPTX
    success = builder.build_from_json(
        args.input,
        args.output
    )

    if success:
        print(f"\n完成！PPTX 文件已保存到: {args.output}")
    else:
        print(f"\n失败：无法生成 PPTX 文件")


if __name__ == '__main__':
    main()
