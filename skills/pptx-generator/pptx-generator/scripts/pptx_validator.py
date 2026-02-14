#!/usr/bin/env python3
"""
PPTX Validator - Validate generated .pptx files.

This module validates .pptx files to ensure they are
correctly generated and can be opened in PowerPoint.
"""

import argparse
import sys
from pathlib import Path
from typing import Optional

try:
    from pptx import Presentation
except ImportError:
    print("错误: 未安装 python-pptx")
    print("请运行: pip install python-pptx")
    sys.exit(1)


# =============================================================================
# PPTX Validator
# =============================================================================

class PPTXValidator:
    """Validator for .pptx files."""

    def __init__(self):
        """Initialize PPTX validator."""
        self.errors = []
        self.warnings = []

    def validate(self, pptx_path: str) -> bool:
        """
        Validate .pptx file.

        Args:
            pptx_path: Path to .pptx file.

        Returns:
            True if valid, False otherwise.
        """
        print(f"验证 PPTX 文件: {pptx_path}\n")

        # Check file exists
        if not Path(pptx_path).exists():
            self.errors.append(f"文件不存在: {pptx_path}")
            self._print_results()
            return False

        # Check file extension
        if not pptx_path.lower().endswith('.pptx'):
            self.warnings.append("文件扩展名不是 .pptx")

        try:
            # Open presentation
            presentation = Presentation(pptx_path)
        except Exception as e:
            self.errors.append(f"无法打开 PPTX 文件: {e}")
            self._print_results()
            return False

        # Validate slides
        self._validate_slides(presentation)

        # Print results
        self._print_results()

        return len(self.errors) == 0

    def _validate_slides(self, presentation) -> None:
        """
        Validate slides in presentation.

        Args:
            presentation: Presentation object.
        """
        slide_count = len(presentation.slides)

        if slide_count == 0:
            self.errors.append("演示文稿中没有幻灯片")
        else:
            print(f"✓ 幻灯片数量: {slide_count}\n")

        # Validate each slide
        for i, slide in enumerate(presentation.slides, 1):
            self._validate_slide(slide, i)

    def _validate_slide(self, slide, index: int) -> None:
        """
        Validate a slide.

        Args:
            slide: Slide object.
            index: Slide index.
        """
        shape_count = len(slide.shapes)

        if shape_count == 0:
            self.warnings.append(f"幻灯片 {index} 没有内容")
        else:
            # Check for title
            has_title = False
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text and shape == slide.shapes.title:
                            has_title = True
                            break

            if not has_title:
                self.warnings.append(f"幻灯片 {index} 没有标题")

    def _print_results(self) -> None:
        """Print validation results."""
        if self.errors:
            print(f"❌ 发现 {len(self.errors)} 个错误：\n")
            for i, error in enumerate(self.errors, 1):
                print(f"  {i}. {error}")
            print()

        if self.warnings:
            print(f"⚠️  发现 {len(self.warnings)} 个警告：\n")
            for i, warning in enumerate(self.warnings, 1):
                print(f"  {i}. {warning}")
            print()

        if not self.errors and not self.warnings:
            print("✓ PPTX 文件验证通过！\n")


# =============================================================================
# Main Function
# =============================================================================

def main():
    """Main entry point."""
    parser = argparse.ArgumentParser(
        description='Validate generated .pptx files'
    )
    parser.add_argument(
        '--input', '-i',
        required=True,
        help='Input .pptx file path'
    )

    args = parser.parse_args()

    # Create validator
    validator = PPTXValidator()

    # Validate
    is_valid = validator.validate(args.input)

    # Exit with appropriate code
    sys.exit(0 if is_valid else 1)


if __name__ == '__main__':
    main()
