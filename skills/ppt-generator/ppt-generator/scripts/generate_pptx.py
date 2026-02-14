#!/usr/bin/env python3
"""
PPT 生成脚本

功能：根据 JSON 格式的数据生成 PowerPoint (.pptx) 文件
依赖：python-pptx>=0.6.21

使用方法：
    python generate_pptx.py --input ./ppt_data.json --output ./presentation.pptx

输入格式：详见 references/ppt_structure_guide.md
"""

import json
import argparse
import sys
from pathlib import Path
from typing import Dict, List, Any

try:
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.enum.text import PP_ALIGN
except ImportError as e:
    print(f"错误：缺少依赖库 {e}")
    print("请安装依赖：pip install python-pptx>=0.6.21")
    sys.exit(1)


class PPTGenerator:
    """PPT 生成器类"""

    # 布局类型映射表
    LAYOUT_MAP = {
        'TitleSlide': 0,           # 标题页
        'TitleAndContent': 1,      # 标题和内容
        'TwoColumnText': 1,        # 双栏文本（使用 TitleAndContent 布局手动分栏）
        'SectionHeader': 2,        # 章节标题页
        'ContentWithCaption': 1,   # 带说明的内容
        'BulletList': 1,           # 项目符号列表
        'BlankSlide': 6,           # 空白页
    }

    def __init__(self, input_json: Dict[str, Any]):
        """
        初始化生成器

        Args:
            input_json: 解析后的 JSON 数据
        """
        self.data = input_json
        self.prs = None
        self._validate_input()

    def _validate_input(self):
        """验证输入数据的合法性"""
        if 'metadata' not in self.data:
            raise ValueError("缺少 metadata 字段")

        if 'slides' not in self.data or not isinstance(self.data['slides'], list):
            raise ValueError("缺少 slides 字段或格式错误")

        if len(self.data['slides']) == 0:
            raise ValueError("slides 不能为空")

        # 验证每个幻灯片
        for idx, slide in enumerate(self.data['slides']):
            if 'layout' not in slide:
                raise ValueError(f"幻灯片 {idx + 1} 缺少 layout 字段")

            if slide['layout'] not in self.LAYOUT_MAP:
                valid_layouts = ', '.join(self.LAYOUT_MAP.keys())
                raise ValueError(
                    f"幻灯片 {idx + 1} 的 layout 值 '{slide['layout']}' 无效。"
                    f"有效值：{valid_layouts}"
                )

            if 'title' not in slide or not slide['title']:
                raise ValueError(f"幻灯片 {idx + 1} 缺少 title 字段或为空")

            if 'content' not in slide or not isinstance(slide['content'], list):
                raise ValueError(f"幻灯片 {idx + 1} 的 content 字段必须是数组")

    def generate(self) -> Presentation:
        """
        生成 PPT 对象

        Returns:
            Presentation: 生成的 PPT 对象
        """
        # 创建演示文稿对象
        self.prs = Presentation()

        # 设置元数据
        self._set_metadata()

        # 生成所有幻灯片
        for slide_data in self.data['slides']:
            self._add_slide(slide_data)

        return self.prs

    def _set_metadata(self):
        """设置演示文稿元数据"""
        metadata = self.data.get('metadata', {})

        # 设置核心属性
        if 'title' in metadata:
            self.prs.core_properties.title = metadata['title']

        if 'author' in metadata:
            self.prs.core_properties.author = metadata['author']

        if 'subject' in metadata:
            self.prs.core_properties.subject = metadata['subject']

        if 'keywords' in metadata:
            self.prs.core_properties.keywords = metadata['keywords']

    def _add_slide(self, slide_data: Dict[str, Any]):
        """
        添加一张幻灯片

        Args:
            slide_data: 幻灯片数据字典
        """
        layout_type = slide_data['layout']
        title = slide_data['title']
        content = slide_data.get('content', [])
        notes = slide_data.get('notes', '')

        # 获取布局索引
        layout_idx = self.LAYOUT_MAP[layout_type]
        slide_layout = self.prs.slide_layouts[layout_idx]

        # 添加幻灯片
        slide = self.prs.slides.add_slide(slide_layout)

        # 设置标题
        if layout_type != 'BlankSlide':
            self._set_title(slide, title)

        # 设置内容
        if layout_type == 'TwoColumnText':
            self._set_two_column_content(slide, content)
        elif layout_type == 'TitleSlide':
            self._set_title_slide_content(slide, content)
        elif content:
            self._set_content(slide, content, layout_type)

        # 设置备注
        if notes:
            self._set_notes(slide, notes)

    def _set_title(self, slide, title_text: str):
        """
        设置幻灯片标题

        Args:
            slide: 幻灯片对象
            title_text: 标题文本
        """
        # 尝试查找标题占位符
        for shape in slide.placeholders:
            if 'Title' in shape.name:
                text_frame = shape.text_frame
                text_frame.text = title_text

                # 设置标题样式
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(36)
                    paragraph.font.bold = True
                break

    def _set_content(self, slide, content: List[str], layout_type: str):
        """
        设置幻灯片内容

        Args:
            slide: 幻灯片对象
            content: 内容列表
            layout_type: 布局类型
        """
        # 查找内容占位符
        content_placeholder = None
        for shape in slide.placeholders:
            if 'Content' in shape.name or 'Body' in shape.name:
                content_placeholder = shape
                break

        if content_placeholder is None:
            # 如果没有占位符，尝试添加文本框
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4)
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
        else:
            text_frame = content_placeholder.text_frame

        # 清空现有内容
        text_frame.clear()

        # 添加内容
        for idx, item in enumerate(content):
            if idx == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()

            paragraph.text = str(item)
            paragraph.font.size = Pt(20)
            paragraph.space_after = Pt(10)

            # 对于 BulletList 设置项目符号
            if layout_type == 'BulletList':
                paragraph.level = 0

    def _set_two_column_content(self, slide, content: List[str]):
        """
        设置双栏内容

        Args:
            slide: 幻灯片对象
            content: 内容列表
        """
        # 解析左右栏
        left_content = []
        right_content = []
        current_side = 'left'

        for item in content:
            item_str = str(item).strip()

            # 检测标记
            if '【右】' in item_str:
                current_side = 'right'
                cleaned_item = item_str.replace('【右】', '').strip()
                if cleaned_item:
                    right_content.append(cleaned_item)
            elif '【左】' in item_str:
                current_side = 'left'
                cleaned_item = item_str.replace('【左】', '').strip()
                if cleaned_item:
                    left_content.append(cleaned_item)
            elif item_str.startswith('[') and ']' in item_str:
                # 可能的格式：[左栏] 或 [右栏]
                if '右' in item_str:
                    current_side = 'right'
                    cleaned_item = item_str.split(']', 1)[1].strip()
                    if cleaned_item:
                        right_content.append(cleaned_item)
                else:
                    current_side = 'left'
                    cleaned_item = item_str.split(']', 1)[1].strip()
                    if cleaned_item:
                        left_content.append(cleaned_item)
            else:
                # 无标记，使用当前栏
                if current_side == 'left':
                    left_content.append(item_str)
                else:
                    right_content.append(item_str)

        # 手动添加两个文本框实现双栏
        # 左栏
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(4))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True

        for idx, item in enumerate(left_content):
            if idx == 0:
                paragraph = left_frame.paragraphs[0]
            else:
                paragraph = left_frame.add_paragraph()

            paragraph.text = item
            paragraph.font.size = Pt(18)
            paragraph.space_after = Pt(8)

        # 右栏
        right_box = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4), Inches(4))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True

        for idx, item in enumerate(right_content):
            if idx == 0:
                paragraph = right_frame.paragraphs[0]
            else:
                paragraph = right_frame.add_paragraph()

            paragraph.text = item
            paragraph.font.size = Pt(18)
            paragraph.space_after = Pt(8)

    def _set_title_slide_content(self, slide, content: List[str]):
        """
        设置标题页内容

        Args:
            slide: 幻灯片对象
            content: 内容列表
        """
        if not content:
            return

        # 查找副标题占位符
        for shape in slide.placeholders:
            if 'Subtitle' in shape.name or 'Date' in shape.name or 'Text Placeholder' in shape.name:
                text_frame = shape.text_frame
                text_frame.text = '\n'.join(content)

                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(24)
                    paragraph.alignment = PP_ALIGN.CENTER
                break

    def _set_notes(self, slide, notes_text: str):
        """
        设置备注

        Args:
            slide: 幻灯片对象
            notes_text: 备注文本
        """
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text

    def save(self, output_path: str):
        """
        保存 PPT 文件

        Args:
            output_path: 输出文件路径
        """
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        self.prs.save(str(output_path))
        print(f"✓ PPT 文件已生成: {output_path.absolute()}")


def main():
    """主函数"""
    parser = argparse.ArgumentParser(description='根据 JSON 数据生成 PowerPoint 文件')
    parser.add_argument('--input', '-i', required=True, help='输入 JSON 文件路径')
    parser.add_argument('--output', '-o', required=True, help='输出 PPT 文件路径')

    args = parser.parse_args()

    # 读取输入文件
    input_path = Path(args.input)
    if not input_path.exists():
        print(f"错误：输入文件不存在: {input_path}")
        sys.exit(1)

    print(f"正在读取输入文件: {input_path}")

    try:
        with open(input_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
    except json.JSONDecodeError as e:
        print(f"错误：JSON 解析失败 - {e}")
        sys.exit(1)
    except Exception as e:
        print(f"错误：读取文件失败 - {e}")
        sys.exit(1)

    # 生成 PPT
    print("正在生成 PPT...")

    try:
        generator = PPTGenerator(data)
        generator.generate()
        generator.save(args.output)
        print("✓ 生成成功！")
    except ValueError as e:
        print(f"错误：数据验证失败 - {e}")
        sys.exit(1)
    except Exception as e:
        print(f"错误：生成失败 - {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == '__main__':
    main()
