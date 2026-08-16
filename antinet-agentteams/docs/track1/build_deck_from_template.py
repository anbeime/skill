"""Build Antinet GOAI track1 proposal deck following the official 8-chapter template.

Uses python-pptx, no network. Output: docs/track1/Antinet_GOAI_track1_v2.pptx
Structure mirrors the Datawhale template:
  Cover -> P0 one-pager -> TOC -> Ch1-8 (each divider + content) -> Demo video.
"""

import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "Antinet_GOAI_track1.pptx")

# ---------- palette (matches Datawhale template: navy + orange) ----------
NAVY = RGBColor(26, 31, 58)
NAVY_LIGHT = RGBColor(37, 44, 74)
NAVY_SOFT = RGBColor(55, 65, 105)
ORANGE = RGBColor(255, 107, 53)
ORANGE_DARK = RGBColor(230, 85, 35)
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(250, 251, 253)
CARD_BG = RGBColor(240, 243, 247)
CARD_BG_DARK = RGBColor(45, 54, 92)
TEXT_DARK = RGBColor(26, 31, 58)
TEXT_GREY = RGBColor(90, 95, 115)
GOLD = RGBColor(184, 138, 43)

# slide size 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

FONT = "Microsoft YaHei"


def add_rect(slide, x, y, w, h, fill, line=None, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = shadow
    return shp


def set_run(run, text, size, color, bold=False, italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT
    # east asian font for chinese
    run.font.language_id = 0x0804
    return run


def add_textbox(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
                valign=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    set_run(p.add_run(), text, size, color, bold, italic)
    return box


def add_bullets(slide, x, y, w, h, items, size=15, color=TEXT_DARK, bold_lead=False,
                line_space=1.15, indent=0.15):
    """items: list of strings; inline **bold** supported."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        p.line_spacing = line_space
        # split inline bold
        parts = re.split(r"(\*\*.*?\*\*)", item)
        for part in parts:
            if not part:
                continue
            run = p.add_run()
            if part.startswith("**") and part.endswith("**"):
                run.text = part[2:-2]
                run.font.bold = True
            else:
                run.text = part
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = FONT
            run.font.language_id = 0x0804
    return box


def add_footer(slide, text, dark_bg=False):
    color = RGBColor(160, 165, 185) if dark_bg else TEXT_GREY
    add_textbox(slide, Inches(0.35), Inches(7.05), Inches(12.6), Inches(0.35),
                text, 9, color)


# =========================== slides ===========================

def slide_cover():
    slide = prs.slides.add_slide(BLANK)
    # background
    add_rect(slide, 0, 0, SW, SH, NAVY)
    # decorative circles
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1.2), Inches(4.5), Inches(4.5))
    c1.fill.solid(); c1.fill.fore_color.rgb = NAVY_LIGHT; c1.line.fill.background()
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(5.0), Inches(3.5), Inches(3.5))
    c2.fill.solid(); c2.fill.fore_color.rgb = NAVY_LIGHT; c2.line.fill.background()
    c3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.5), Inches(5.5), Inches(2.8), Inches(2.8))
    c3.fill.solid(); c3.fill.fore_color.rgb = NAVY_SOFT; c3.line.fill.background()

    add_textbox(slide, Inches(0.7), Inches(2.0), Inches(10), Inches(0.5),
                "GOAI 世界人工智能开源大赛", 16, ORANGE, bold=True)
    add_textbox(slide, Inches(0.7), Inches(2.55), Inches(10), Inches(1.0),
                "Antinet·八官署", 54, WHITE, bold=True)
    add_textbox(slide, Inches(0.7), Inches(3.65), Inches(10), Inches(0.55),
                "多 AI 知识库中间层", 28, WHITE)
    add_textbox(slide, Inches(0.7), Inches(4.25), Inches(11), Inches(0.5),
                "软件研发全流程协同 · Agent Infra 新智基座 · 方向三", 18, RGBColor(180, 185, 205))
    add_footer(slide, "Antinet·八官署 — 面向复杂知识任务的多智能体基础设施", dark_bg=True)


def slide_p0():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SW, SH, WHITE)
    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(3), Inches(0.4),
                "P0 · 一页纸速览", 14, ORANGE, bold=True)
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(4), Inches(0.6),
                "作品简介", 32, TEXT_DARK, bold=True)

    cards = [
        ("项目名称", "Antinet·八官署：多 AI 知识库中间层"),
        ("问题与场景", "研发团队在缺陷/需求/日志/反馈中识别根因并产出可验证修复方案；传统单 Agent 把聚合/定位/修复/验证混在一起，难审计、难降级、难复现。"),
        ("核心解决方案", "以「明朝内阁」治理体系映射为八官署多智能体框架；四色卡片结构化传递上下文；太史阁先读后写，实现知识复利。"),
        ("创新点与差异化", "① 可溯源四色卡片 ② AgentTeams 协同基点 ③ 15 热插拔 Skill ④ 记忆+共享状态+可观测 ⑤ 熔断降级不崩盘"),
        ("开放 / 复用价值", "Apache 2.0 开源；能力可迁移至企业 AI 编码合规审计、跨工具创意生产、智能客服知识自演化、个性化学习教育陪练。"),
        ("当前进展", "八署串连通跑，verify_production.py PASS；已接入真实本地 NPU（Genie qwen2.5vl3b）；Agent Identity / Skill 体系 / AgentTeams 映射文档齐备。"),
    ]
    cols = [0.5, 4.55, 8.6]
    rows = [1.65, 4.25]
    w, h = Inches(3.75), Inches(2.35)
    for idx, (title, body) in enumerate(cards):
        c = idx % 3
        r = idx // 3
        x = Inches(cols[c])
        y = Inches(rows[r])
        card = add_rect(slide, x, y, w, h, CARD_BG, shadow=False)
        # orange dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.18), y + Inches(0.18), Inches(0.14), Inches(0.14))
        dot.fill.solid(); dot.fill.fore_color.rgb = ORANGE; dot.line.fill.background()
        add_textbox(slide, x + Inches(0.42), y + Inches(0.14), w - Inches(0.55), Inches(0.35),
                    title, 15, TEXT_DARK, bold=True)
        add_textbox(slide, x + Inches(0.18), y + Inches(0.55), w - Inches(0.36), h - Inches(0.7),
                    body, 12, TEXT_GREY)
    add_footer(slide, "Antinet·八官署 — GOAI 赛道一 Agent Infra 新智基座")


def slide_toc():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SW, SH, NAVY)
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(-1.0), Inches(4.2), Inches(4.2))
    c1.fill.solid(); c1.fill.fore_color.rgb = NAVY_LIGHT; c1.line.fill.background()
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.2), Inches(5.2), Inches(3.2), Inches(3.2))
    c2.fill.solid(); c2.fill.fore_color.rgb = NAVY_LIGHT; c2.line.fill.background()

    add_textbox(slide, Inches(0.7), Inches(0.7), Inches(3), Inches(0.6),
                "目录", 36, WHITE, bold=True)

    chapters = [
        ("1", "场景与价值"),
        ("2", "方案总览"),
        ("3", "多 Agent 协同设计"),
        ("4", "Skills 工具体系"),
        ("5", "工程落地、运行验证与安全可审计"),
        ("6", "开源开放计划"),
        ("7", "落地计划与进展"),
        ("8", "Demo 视频"),
    ]
    xs = [0.7, 6.9]
    y0 = 1.7
    h = 0.72
    for i, (num, title) in enumerate(chapters):
        c = i % 2
        r = i // 2
        x = Inches(xs[c])
        y = Inches(y0 + r * (h + 0.22))
        w = Inches(5.6)
        card = add_rect(slide, x, y, w, Inches(h), CARD_BG_DARK, shadow=False)
        # orange badge
        badge = add_rect(slide, x + Inches(0.12), y + Inches(0.12), Inches(0.48), Inches(0.48), ORANGE, shadow=False)
        add_textbox(slide, x + Inches(0.12), y + Inches(0.12), Inches(0.48), Inches(0.48),
                    num, 16, WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x + Inches(0.78), y + Inches(0.12), Inches(4.6), Inches(0.48),
                    title, 16, WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide, "八章结构对齐 Agent Infra 初赛评分维度", dark_bg=True)


def slide_divider(chapter_no, chapter_title, dimension, weight):
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SW, SH, WHITE)
    # chapter badge
    add_textbox(slide, Inches(0.6), Inches(2.2), Inches(3), Inches(0.4),
                f"第{chapter_no}章", 14, ORANGE, bold=True)
    add_textbox(slide, Inches(0.6), Inches(2.7), Inches(8), Inches(0.9),
                chapter_title, 44, TEXT_DARK, bold=True)
    # scoring badge top-right
    add_rect(slide, Inches(9.2), Inches(0.6), Inches(3.5), Inches(0.85), RGBColor(255, 245, 240), line=ORANGE)
    add_textbox(slide, Inches(9.35), Inches(0.68), Inches(1.8), Inches(0.3),
                "对应评分维度", 10, ORANGE, bold=True)
    add_textbox(slide, Inches(9.35), Inches(0.98), Inches(2.4), Inches(0.4),
                dimension, 12, TEXT_DARK, bold=True)
    add_rect(slide, Inches(11.9), Inches(0.6), Inches(0.8), Inches(0.85), ORANGE, shadow=False)
    add_textbox(slide, Inches(11.9), Inches(0.6), Inches(0.8), Inches(0.85),
                weight, 18, WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide, "Antinet·八官署 — GOAI 赛道一 Agent Infra 新智基座")


def slide_ch1_scene():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SW, SH, WHITE)
    add_textbox(slide, Inches(0.55), Inches(0.35), Inches(1.2), Inches(0.32),
                "第一章", 12, ORANGE, bold=True)
    add_textbox(slide, Inches(0.55), Inches(0.7), Inches(5), Inches(0.6),
                "场景与价值", 32, TEXT_DARK, bold=True)
    add_textbox(slide, Inches(0.55), Inches(1.35), Inches(11.8), Inches(0.45),
                "目标用户、核心痛点、方向三实例化与行业可复制性", 14, TEXT_GREY)

    bullets = [
        "**目标用户**：开发者与研发团队，需从 Issue / 日志 / 用户反馈中快速定位根因并产出可验证修复。",
        "**真实痛点**：信息分散、根因定位靠经验、修复方案难审计、复盘知识随人走、单 Agent 把多阶段任务混在一起。",
        "**方向三实例化**：缺陷/需求聚合 → 代码根因定位 → 修复生成与执行 → 测试/发布确认 → 上线复盘与知识沉淀。",
        "**行业可复制**：四色卡片作为通用知识格式，同样适用于金融风控、客服复盘、运维复盘、材料研发。",
        "**差异化定位**：不是单 Agent Demo，而是让多个 AI 围绕同一份结构化知识协同的「知识库中间层」。",
    ]
    add_bullets(slide, Inches(0.6), Inches(2.0), Inches(6.8), Inches(5.2), bullets, size=15)

    # 4-scene cards on right
    scenes = [
        ("企业 AI 编码合规审计", "Secrets 出域防控、改动审计"),
        ("跨工具创意生产管线", "素材/版本/决策集中管理"),
        ("智能客服知识自演化", "从静态文档到带溯源自活体"),
        ("个性化学习教育陪练", "错点追踪与可解释学习路径"),
    ]
    x0, y0 = 7.5, 2.0
    wf, hf = 2.65, 1.25
    gap = 0.22
    for i, (title, desc) in enumerate(scenes):
        x = Inches(x0 + (i % 2) * (wf + gap))
        y = Inches(y0 + (i // 2) * (hf + gap))
        w = Inches(wf)
        h = Inches(hf)
        card = add_rect(slide, x, y, w, h, CARD_BG, shadow=False)
        add_textbox(slide, x + Inches(0.12), y + Inches(0.12), w - Inches(0.24), Inches(0.45),
                    title, 12, TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.12), y + Inches(0.58), w - Inches(0.24), Inches(0.6),
                    desc, 10, TEXT_GREY, align=PP_ALIGN.CENTER)
    add_footer(slide, "四个真实付费场景证明底座可跨行业复用")


def slide_ch2_overview():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SW, SH, WHITE)
    add_textbox(slide, Inches(0.55), Inches(0.35), Inches(1.6), Inches(0.32),
                "第二章 · 承上启下", 12, ORANGE, bold=True)
    add_textbox(slide, Inches(0.55), Inches(0.7), Inches(4), Inches(0.6),
                "方案总览", 32, TEXT_DARK, bold=True)
    add_textbox(slide, Inches(0.55), Inches(1.35), Inches(11.8), Inches(0.45),
                "端到端架构与主流程", 14, TEXT_GREY)

    # architecture diagram: stacked layers on the left
    layers = [
        ("用户 / 客户入口", "缺陷、需求、告警、日志、反馈"),
        ("指挥使 · 编排层", "意图识别、任务拆解、异常熔断"),
        ("锦衣卫 · 安检层", "合规扫描、防假全文、密钥检测"),
        ("密卷房 · 解析层", "多格式解析与 OCR 三级 fallback"),
        ("通政司 · 事实层", "事实抽取 → 蓝卡"),
        ("监察院 · 审计层", "逻辑审查、Gap → 绿卡"),
        ("丞相府 · 策略层", "构效假说与行动建议 → 红卡"),
        ("军机处 · 执行层", "核验、产物落盘、报告生成"),
        ("太史阁 · 记忆层", "长期记忆、检索、知识回流、provenance"),
    ]
    x = Inches(0.6)
    y0 = 2.0
    wf, hf = 4.6, 0.50
    gap = 0.07
    for i, (title, desc) in enumerate(layers):
        y = Inches(y0 + i * (hf + gap))
        w = Inches(wf)
        h = Inches(hf)
        fill = NAVY if i == 0 or i == len(layers) - 1 else CARD_BG
        txt = WHITE if i == 0 or i == len(layers) - 1 else TEXT_DARK
        add_rect(slide, x, y, w, h, fill, shadow=False)
        add_textbox(slide, x + Inches(0.1), y, Inches(1.6), h,
                    title, 11, txt, bold=True, valign=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x + Inches(1.75), y, Inches(2.75), h,
                    desc, 10, txt, valign=MSO_ANCHOR.MIDDLE)

    # right column bullets
    bullets = [
        "**职责隔离**：8 个具名 Agent 各司其职，每个环节可独立验证、独立降级。",
        "**四色卡片**：蓝事实 / 绿解释 / 黄风险 / 红行动，作为多 AI 间唯一结构化上下文介质。",
        "**记忆总线**：太史阁「先读库再干活、后写库做沉淀」，每次任务都在复用与增值知识。",
        "**设计基点**：以 AgentTeams 角色模型为协同基点，后续迁移只需协议适配。",
    ]
    add_bullets(slide, Inches(5.6), Inches(2.0), Inches(7.0), Inches(5.2), bullets, size=15)
    add_footer(slide, "八署 + 四色卡片 + 记忆总线 = 可审计的多 AI 知识中间层")


def slide_ch3_agents():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SW, SH, WHITE)
    add_textbox(slide, Inches(0.55), Inches(0.35), Inches(1.2), Inches(0.32),
                "第三章", 12, ORANGE, bold=True)
    add_textbox(slide, Inches(0.55), Inches(0.7), Inches(7), Inches(0.6),
                "多 Agent 协同设计", 32, TEXT_DARK, bold=True)
    add_textbox(slide, Inches(0.55), Inches(1.35), Inches(11.8), Inches(0.45),
                "Agent 分工、任务拆解、上下文传递、异常处理与安全边界", 14, TEXT_GREY)

    # pipeline boxes across top
    agents = [
        ("指挥使", "编排"),
        ("锦衣卫", "安检"),
        ("密卷房", "解析"),
        ("通政司", "蓝卡"),
        ("监察院", "绿卡"),
        ("丞相府", "红卡"),
        ("军机处", "核验"),
        ("太史阁", "记忆"),
    ]
    x0 = 0.55
    y0 = 1.95
    wf, hf = 1.35, 0.85
    gap = 0.16
    for i, (name, role) in enumerate(agents):
        x = Inches(x0 + i * (wf + gap))
        y = Inches(y0)
        w = Inches(wf)
        h = Inches(hf)
        card = add_rect(slide, x, y, w, h, CARD_BG, shadow=False)
        add_textbox(slide, x, y + Inches(0.08), w, Inches(0.35),
                    name, 13, TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(0.45), w, Inches(0.35),
                    role, 11, ORANGE, align=PP_ALIGN.CENTER)
        # arrows between
        if i < len(agents) - 1:
            ax = x + w
            add_textbox(slide, ax, y + Inches(0.28), Inches(gap), Inches(0.25),
                        "→", 18, TEXT_GREY, align=PP_ALIGN.CENTER)

    bullets = [
        "**角色编排**：8 个具名 Agent，身份、能力边界与协同关系见《Agent Identity 清单》（附录 A）。",
        "**任务拆解**：指挥使状态机将用户主题拆为检索 / 解析 / 抽取 / 审查 / 假说 / 核验 / 回流子任务。",
        "**上下文传递**：四色卡片 cite 链实现传递；每署只接收上游卡片，不读原始噪声。",
        "**异常与冲突**：任一官署失败即触发指挥使熔断与规则引擎降级，大模型掉线/解析失败不崩盘。",
        "**安全边界**：锦衣卫是唯一出域审批方；高风险动作（发布/合并）由指挥使状态机执行审批 / 回滚 / 审计。",
    ]
    add_bullets(slide, Inches(0.6), Inches(3.1), Inches(12.0), Inches(4.0), bullets, size=15)
    add_footer(slide, "多 Agent 协同以职责隔离与状态可见为核心")


def slide_ch4_skills():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SW, SH, WHITE)
    add_textbox(slide, Inches(0.55), Inches(0.35), Inches(2.2), Inches(0.32),
                "第四章 · 本赛题必选项", 12, ORANGE, bold=True)
    add_textbox(slide, Inches(0.55), Inches(0.7), Inches(5), Inches(0.6),
                "Skill 工程体系", 32, TEXT_DARK, bold=True)
    add_textbox(slide, Inches(0.55), Inches(1.35), Inches(11.8), Inches(0.45),
                "15 个真实可运行 Skill 与完整规格", 14, TEXT_GREY)

    # skill categories as 4 cards
    cats = [
        ("抽取类", "four_color_cards\nmarkdown_converter\nlocal_audio_processor"),
        ("知识类", "knowledge_graph\nbook_skill\nview_manager"),
        ("展示类", "infographic\nhtml_report\nreport_automation\nppt_structure_draft\nchart_recommendation"),
        ("治理类", "card_filter\nmarkdown_formatter\ninvoice_skill"),
    ]
    x0 = 0.6
    y0 = 1.95
    wf, hf = 2.85, 1.7
    gap = 0.28
    for i, (title, body) in enumerate(cats):
        x = Inches(x0 + i * (wf + gap))
        y = Inches(y0)
        w = Inches(wf)
        h = Inches(hf)
        card = add_rect(slide, x, y, w, h, CARD_BG, shadow=False)
        add_textbox(slide, x + Inches(0.12), y + Inches(0.12), w - Inches(0.24), Inches(0.35),
                    title, 14, TEXT_DARK, bold=True)
        add_textbox(slide, x + Inches(0.12), y + Inches(0.52), w - Inches(0.24), h - Inches(0.65),
                    body, 11, TEXT_GREY)

    bullets = [
        "**单个 Skill 规格**：名称 / 用途 / 输入输出 / 调用条件 / 依赖工具 / 失败处理 / 安全边界 / 复用价值 / 与协同流程的关系。",
        "**复用与生命周期**：统一 Skill 基类（name/description/category/parameters_schema/execute），支持热插拔、版本化、失败降级。",
        "**生态对接**：核心抽象层与 github.com/anbeime/skill 公开技能商店接口对齐，迁移到新运行时只需协议适配。",
        "**赛道要求**：Skill 是必选项；本方案 15 个真实 Skill 可独立部署，亦可被外部 Agent 复用。",
    ]
    add_bullets(slide, Inches(0.6), Inches(3.85), Inches(12.0), Inches(3.3), bullets, size=14)
    add_footer(slide, "Skill 是能力抽象层，而非一次性行为")


def slide_ch5_engineering():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SW, SH, WHITE)
    add_textbox(slide, Inches(0.55), Inches(0.35), Inches(1.2), Inches(0.32),
                "第五章", 12, ORANGE, bold=True)
    add_textbox(slide, Inches(0.55), Inches(0.7), Inches(9), Inches(0.6),
                "工程落地、运行验证与安全可审计", 32, TEXT_DARK, bold=True)
    add_textbox(slide, Inches(0.55), Inches(1.35), Inches(11.8), Inches(0.45),
                "可运行性、运行证据、可观测与安全治理", 14, TEXT_GREY)

    # 2x2 evidence cards
    cards = [
        ("可运行性", "python scripts/run_survey.py 主链路跑通\nverify_production.py 输出 consistency:PASS，退出码 0\n可接入 CI 幂等重跑"),
        ("运行证据", "verify_report.json：parsed_oa_ratio 7/10\nprovenance/knowledge.md 全链路回放\n四色卡片样例可审计"),
        ("可观测 / 检索", "太史阁 = 记忆存储 + 共享状态管理\nprovenance 日志 = Trace / Log / Metrics\n超额满足 2/4 项要求"),
        ("安全治理", "默认本地 NPU，数据不出域\n锦衣卫合规扫描 + 密钥检测\n指挥使状态机审批 / 回滚"),
    ]
    x0, y0 = 0.6, 2.0
    wf, hf = 5.9, 1.35
    gapx, gapy = 0.3, 0.25
    for i, (title, body) in enumerate(cards):
        x = Inches(x0 + (i % 2) * (wf + gapx))
        y = Inches(y0 + (i // 2) * (hf + gapy))
        w = Inches(wf)
        h = Inches(hf)
        card = add_rect(slide, x, y, w, h, CARD_BG, shadow=False)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.12), w - Inches(0.3), Inches(0.35),
                    title, 15, TEXT_DARK, bold=True)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.5), w - Inches(0.3), h - Inches(0.65),
                    body, 12, TEXT_GREY)

    add_textbox(slide, Inches(0.6), Inches(5.1), Inches(12.0), Inches(0.35),
                "云产品选型边界：Nacos / Higress / PolarDB / RocketMQ / LoongSuite 均可等价替换，迁移成本已在映射文档中逐条论证。",
                13, TEXT_DARK, bold=True)
    add_footer(slide, "诚实披露现状：初赛为方案设计，完整代码仓库接入将于复赛落地")


def slide_ch6_opensource():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SW, SH, WHITE)
    add_textbox(slide, Inches(0.55), Inches(0.35), Inches(1.2), Inches(0.32),
                "第六章", 12, ORANGE, bold=True)
    add_textbox(slide, Inches(0.55), Inches(0.7), Inches(5), Inches(0.6),
                "开放 / 开源计划", 32, TEXT_DARK, bold=True)
    add_textbox(slide, Inches(0.55), Inches(1.35), Inches(11.8), Inches(0.45),
                "可复用成果、接口契约与开源协议", 14, TEXT_GREY)

    bullets = [
        "**开源协议**：Apache 2.0；本方案 materials-agent/ 仓库结构公开。",
        "**可复用成果**：四色卡片开放格式、八署编排运行时、15 Skill 抽象层、provenance 日志格式。",
        "**接口契约**：等价外部工具集成契约（协议 / 鉴权 / Schema / 错误 / 审计 / 迁移成本），不使用 MCP 也给出等价方案。",
        "**文档与示例**：Agent Identity 清单、AgentTeams 映射文档、Skill 工程体系说明、DEPENDENCIES.md 第三方依赖边界。",
        "**数据集**：OpenAlex（CC0）合法复用；闭源 API 的接口与降级逻辑均开源。",
        "**Skill 生态**：核心 Skill 与 github.com/anbeime/skill 公开技能商店对齐，可被外部 Agent 直接复用。",
    ]
    add_bullets(slide, Inches(0.6), Inches(2.0), Inches(7.6), Inches(4.6), bullets, size=15)

    # right side callout box
    add_rect(slide, Inches(8.7), Inches(2.0), Inches(3.9), Inches(3.0), CARD_BG, shadow=False)
    add_textbox(slide, Inches(8.85), Inches(2.15), Inches(3.6), Inches(0.4),
                "开源目标", 14, ORANGE, bold=True)
    add_textbox(slide, Inches(8.85), Inches(2.6), Inches(3.6), Inches(2.25),
                "把「多 AI 可审计记忆总线」做成标准件，而不是绑定在某个封闭产品里。",
                13, TEXT_DARK)
    add_footer(slide, "开放协议与可替换性是方案的核心壁垒")


def slide_ch7_roadmap():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SW, SH, WHITE)
    add_textbox(slide, Inches(0.55), Inches(0.35), Inches(2.6), Inches(0.32),
                "第七章 · 对应「当前进展」", 12, ORANGE, bold=True)
    add_textbox(slide, Inches(0.55), Inches(0.7), Inches(5), Inches(0.6),
                "落地计划与进展", 32, TEXT_DARK, bold=True)
    add_textbox(slide, Inches(0.55), Inches(1.35), Inches(11.8), Inches(0.45),
                "当前进展、里程碑与风险控制", 14, TEXT_GREY)

    # timeline
    milestones = [
        ("当前", "八署串连通跑\nLLM 真实接入\n15 Skill 就绪"),
        ("初赛 8/16", "方案 PPT + 作品简介\n附件 ZIP 提交"),
        ("复赛 8/25", "可执行 AgentTeams 代码包\n+ 完整 Demo"),
        ("决赛 9/22", "杭州现场答辩\n展示可运行闭环"),
    ]
    x0 = 0.9
    y0 = 2.2
    wf, hf = 2.65, 1.9
    gap = 0.35
    for i, (title, body) in enumerate(milestones):
        x = Inches(x0 + i * (wf + gap))
        y = Inches(y0)
        w = Inches(wf)
        h = Inches(hf)
        card = add_rect(slide, x, y, w, h, CARD_BG if i > 0 else NAVY, shadow=False)
        txt = WHITE if i == 0 else TEXT_DARK
        add_textbox(slide, x + Inches(0.15), y + Inches(0.15), w - Inches(0.3), Inches(0.4),
                    title, 15, ORANGE, bold=True)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.62), w - Inches(0.3), Inches(1.15),
                    body, 12, txt)
        if i < len(milestones) - 1:
            add_textbox(slide, x + w, y + h / 2 - Inches(0.15), Inches(gap), Inches(0.3),
                        "→", 22, ORANGE, align=PP_ALIGN.CENTER)

    bullets = [
        "**当前进展**：八署 pipeline 可复现、verify PASS、真实本地 NPU 接入、Agent Identity / Skill / AgentTeams 映射文档齐备。",
        "**复赛目标**：迁移八署至 AgentTeams 运行时，接入 Higress / Nacos / PolarDB，完成代码仓库 PR/测试/发布完整闭环。",
        "**风险控制**：初赛阶段诚实披露降级现状（解析预存全文模拟 7/10、代码仓库工具链复赛落地），避免粉饰。",
    ]
    add_bullets(slide, Inches(0.6), Inches(4.35), Inches(12.0), Inches(2.7), bullets, size=14)
    add_footer(slide, "从可运行方案到可执行代码包的演进路径清晰")


def slide_ch8_team():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SW, SH, WHITE)
    add_textbox(slide, Inches(0.55), Inches(0.35), Inches(1.2), Inches(0.32),
                "第八章", 12, ORANGE, bold=True)
    add_textbox(slide, Inches(0.55), Inches(0.7), Inches(5), Inches(0.6),
                "团队介绍", 32, TEXT_DARK, bold=True)
    add_textbox(slide, Inches(0.55), Inches(1.35), Inches(11.8), Inches(0.45),
                "成员背景、分工与成果（请补充真实信息）", 14, TEXT_GREY)

    # placeholder cards for 3 members
    members = [
        ("负责人 / 架构", "多 Agent 系统设计、八署编排、AgentTeams 映射"),
        ("算法 / Skill", "四色卡片模型、Skill 抽象层、本地 NPU 接入"),
        ("工程 / 落地", "可运行链路、验证脚本、provenance 与可观测"),
    ]
    x0 = 0.8
    y0 = 2.2
    wf, hf = 3.7, 2.0
    gap = 0.35
    for i, (role, scope) in enumerate(members):
        x = Inches(x0 + i * (wf + gap))
        y = Inches(y0)
        w = Inches(wf)
        h = Inches(hf)
        card = add_rect(slide, x, y, w, h, CARD_BG, shadow=False)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.2), w - Inches(0.4), Inches(0.4),
                    f"成员 {i+1}", 16, TEXT_DARK, bold=True)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.7), w - Inches(0.4), Inches(0.35),
                    role, 13, ORANGE, bold=True)
        add_textbox(slide, x + Inches(0.2), y + Inches(1.12), w - Inches(0.4), Inches(0.75),
                    scope, 12, TEXT_GREY)

    add_textbox(slide, Inches(0.6), Inches(4.55), Inches(12.0), Inches(0.8),
                "提示：请在此页替换为真实成员信息，包括学校/公司、岗位/专业、核心技能、过往项目/获奖经历及作品合集链接。",
                13, TEXT_GREY, italic=True)
    add_footer(slide, "团队是方案落地的最终执行者")


def slide_demo():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SW, SH, NAVY)
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(-1.0), Inches(4.2), Inches(4.2))
    c1.fill.solid(); c1.fill.fore_color.rgb = NAVY_LIGHT; c1.line.fill.background()

    add_textbox(slide, Inches(0.7), Inches(1.8), Inches(10), Inches(0.5),
                "Demo 视频与提交信息", 40, WHITE, bold=True)
    add_textbox(slide, Inches(0.7), Inches(2.6), Inches(11.5), Inches(0.45),
                "docs/track1/Antinet_GOAI_track1_demo.mp4（45 秒，1080p）", 18, ORANGE, bold=True)

    infos = [
        "**作品名**：Antinet·八官署（软件研发全流程协同 · 多 AI 知识库中间层）",
        "**赛道**：GOAI 赛道一 · Agent Infra 新智基座 · 方向三「软件研发全流程协同」",
        "**开源仓库**：materials-agent/（Apache 2.0）+ github.com/anbeime/skill 公开技能商店",
        "**Demo 入口**：scripts/run_survey.py 主链路 + verify_production.py CI 可复现",
    ]
    add_bullets(slide, Inches(0.7), Inches(3.3), Inches(11.8), Inches(3.0), infos, size=15, color=RGBColor(210, 215, 230))
    add_footer(slide, "八官署证明：多智能体的价值不在参数规模，而在职责隔离与可追溯", dark_bg=True)


# =========================== build ===========================

slide_cover()
slide_p0()
slide_toc()
slide_divider("一", "场景与价值", "场景价值与行业可复制性", "25%")
slide_ch1_scene()
slide_divider("二", "方案总览", "端到端方案与关键技术选型", "—")
slide_ch2_overview()
slide_divider("三", "多 Agent 协同设计", "多 Agent 协同与自主闭环能力", "25%")
slide_ch3_agents()
slide_divider("四", "Skill 工程体系", "Skill 工程体系与生态复用", "25%")
slide_ch4_skills()
slide_divider("五", "工程落地、运行验证与安全可审计", "工程落地与安全可审计", "20%")
slide_ch5_engineering()
slide_divider("六", "开放 / 开源计划", "开放 / 开源贡献", "5%")
slide_ch6_opensource()
slide_divider("七", "落地计划与进展", "当前进展与整体可行性", "—")
slide_ch7_roadmap()
slide_divider("八", "团队介绍", "团队能力与执行可行性", "—")
slide_ch8_team()
slide_demo()

prs.save(OUT)
print("saved:", OUT, "slides:", len(prs.slides))
