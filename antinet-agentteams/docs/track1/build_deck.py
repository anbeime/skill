"""Render docs/track1/deck_track1.md -> Antinet_GOAI_track1.pptx (18 slides).

Parses "## Pn · Title" sections; bullets are "- **bold**：text" or "- text".
Pure python-pptx, no network. Output beside this script.
"""
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(HERE, "deck_track1.md")
OUT = os.path.join(HERE, "Antinet_GOAI_track1.pptx")

INK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0xC8, 0x39, 0x2B)   # 朱红（八署配色）
GOLD = RGBColor(0xB8, 0x8A, 0x2B)
LIGHT = RGBColor(0xF5, 0xF1, 0xE6)    # 宣纸
GREY = RGBColor(0x55, 0x55, 0x55)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def set_text(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Microsoft YaHei"
    return p


def add_para(tf, text, size, color, bold=False, level=0, space=6):
    p = tf.add_paragraph()
    p.level = level
    p.space_after = Pt(space)
    # support inline **bold**
    parts = re.split(r"(\*\*.+?\*\*)", text)
    for part in parts:
        if not part:
            continue
        r = p.add_run()
        if part.startswith("**") and part.endswith("**"):
            r.text = part[2:-2]
            r.font.bold = True
        else:
            r.text = part
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Microsoft YaHei"
    return p


# Parse source
with open(SRC, encoding="utf-8") as f:
    raw = f.read()

slides_md = re.split(r"\n## ", raw)
# slides_md[0] is the header (ignore)
slides = []
for block in slides_md[1:]:
    lines = block.splitlines()
    head = lines[0].strip()
    m = re.match(r"^(P\d+)\s*[·•-]?\s*(.+)$", head)
    if not m:
        continue
    pid, title = m.group(1), m.group(2).strip()
    bullets = []
    for ln in lines[1:]:
        ln = ln.strip()
        if ln.startswith("- "):
            bullets.append(ln[2:].strip())
    slides.append((pid, title, bullets))

for i, (pid, title, bullets) in enumerate(slides):
    slide = prs.slides.add_slide(BLANK)
    # background
    add_rect(slide, 0, 0, SW, SH, LIGHT)
    # top accent bar
    add_rect(slide, 0, 0, SW, Inches(0.18), ACCENT)
    # left spine
    add_rect(slide, 0, 0, Inches(0.18), SH, ACCENT)
    # page tag
    tag = slide.shapes.add_textbox(Inches(0.35), Inches(0.30), Inches(2.2), Inches(0.5))
    set_text(tag.text_frame, f"GOAI · {pid}", 13, GOLD, bold=True)
    # title
    tbox = slide.shapes.add_textbox(Inches(0.35), Inches(0.85), Inches(12.6), Inches(1.0))
    set_text(tbox.text_frame, title, 28, INK, bold=True)
    # divider
    add_rect(slide, Inches(0.35), Inches(1.85), Inches(3.2), Pt(3), GOLD)
    # bullets
    bbox = slide.shapes.add_textbox(Inches(0.45), Inches(2.1), Inches(12.4), Inches(4.9))
    tf = bbox.text_frame
    tf.word_wrap = True
    for idx, b in enumerate(bullets):
        if idx == 0:
            # reuse the empty first paragraph
            p = tf.paragraphs[0]
            c = ACCENT
        else:
            p = tf.add_paragraph()
            c = INK
        p.space_after = Pt(8)
        parts = re.split(r"(\*\*.+?\*\*)", b)
        for part in parts:
            if not part:
                continue
            r = p.add_run()
            if part.startswith("**") and part.endswith("**"):
                r.text = part[2:-2]
                r.font.bold = True
            else:
                r.text = part
            r.font.size = Pt(16)
            r.font.color.rgb = c
            r.font.name = "Microsoft YaHei"
    # footer
    foot = slide.shapes.add_textbox(Inches(0.35), Inches(7.05), Inches(12.6), Inches(0.35))
    set_text(foot.text_frame, "Antinet·八官署 — 面向复杂知识任务的多智能体基础设施  |  GOAI 赛道一 Agent Infra",
             9, GREY)

prs.save(OUT)
print("saved:", OUT, "slides:", len(slides))
